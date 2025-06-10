import io
from typing import List
from uuid import uuid4

from PyPDF2 import PdfReader
import docx
from pptx import Presentation

import nltk
from nltk.tokenize import sent_tokenize
import os


env_chunk_size = int(os.getenv("CHUNK_SIZE", 500))
env_chunk_overlap = int(os.getenv("CHUNK_OVERLAP", 50))


def extract_text_from_pdf(file_bytes: bytes) -> str:
    reader = PdfReader(io.BytesIO(file_bytes))
    texts: List[str] = []
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            texts.append(page_text)
    return "\n".join(texts)


def extract_text_from_docx(file_bytes: bytes) -> str:
    doc = docx.Document(io.BytesIO(file_bytes))
    texts: List[str] = []
    for para in doc.paragraphs:
        texts.append(para.text)
    return "\n".join(texts)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    texts: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)


def extract_text_from_txt(file_bytes: bytes) -> str:
    try:
        return file_bytes.decode("utf-8", errors="ignore")
    except Exception:
        return ""


def chunk_text(text: str, chunk_size: int = None, chunk_overlap: int = None) -> List[dict]:
    """
    Splits `text` into chunks of approximate token length `chunk_size`, overlapping by `chunk_overlap`.
    Uses sentence tokenization to preserve boundaries when possible.
    Returns a list of dicts: [{"chunk_id": UUID, "text": str, "offset": int}, ...]
    """
    if chunk_size is None:
        chunk_size = env_chunk_size
    if chunk_overlap is None:
        chunk_overlap = env_chunk_overlap

    sentences = sent_tokenize(text)
    chunks = []
    current_chunk: List[str] = []
    current_len = 0
    offset = 0

    for sentence in sentences:
        tokens = sentence.split()
        token_count = len(tokens)
        if current_len + token_count > chunk_size:
            chunk_text = " ".join(current_chunk)
            chunks.append({
                "chunk_id": uuid4(),
                "text": chunk_text,
                "offset": offset,
            })
            overlap_tokens = []
            overlap_len = 0
            while current_len > 0 and overlap_len < chunk_overlap:
                last_sentence = current_chunk.pop()
                overlap_tokens.insert(0, last_sentence)
                overlap_len += len(last_sentence.split())
                current_len -= len(last_sentence.split())
            current_chunk = overlap_tokens.copy()
            current_len = overlap_len
            offset += 1
        current_chunk.append(sentence)
        current_len += token_count
    # Final chunk
    if current_chunk:
        chunk_text = " ".join(current_chunk)
        chunks.append({
            "chunk_id": uuid4(),
            "text": chunk_text,
            "offset": offset,
        })
    return chunks
